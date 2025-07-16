import os
from typing import List, Dict, Any, Tuple, Optional
from langchain_groq import ChatGroq
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain.chains import RetrievalQA
from langchain_community.retrievers import BM25Retriever
from langchain.retrievers import EnsembleRetriever
from langchain.schema import Document
import PyPDF2
import pdfplumber
from docx import Document as DocxDocument
import pandas as pd
try:
    import magic
    MAGIC_AVAILABLE = True
except ImportError:
    MAGIC_AVAILABLE = False
    print("Warning: python-magic not available. Using file extension for type detection.")
import tiktoken
import json
import xml.etree.ElementTree as ET
import yaml
import csv
import xlrd
import openpyxl
import pptx
import html
import re
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from datetime import datetime
import base64
from io import BytesIO
import hashlib

# Image processing imports
from PIL import Image, ImageEnhance, ImageFilter
import numpy as np

# Handle EasyOCR with system compatibility
try:
    import easyocr
    EASYOCR_AVAILABLE = True
except ImportError:
    EASYOCR_AVAILABLE = False
    print("Warning: EasyOCR not available. Using Tesseract only for OCR.")
except Exception as e:
    EASYOCR_AVAILABLE = False
    print(f"Warning: EasyOCR initialization failed: {e}. Using Tesseract only for OCR.")

# Handle Tesseract
try:
    import pytesseract
    PYTESSERACT_AVAILABLE = True
except ImportError:
    PYTESSERACT_AVAILABLE = False
    print("Warning: Pytesseract not available. OCR functionality will be limited.")

# Handle OpenCV with fallback
try:
    import cv2
    CV2_AVAILABLE = True
except ImportError:
    CV2_AVAILABLE = False
    print("Warning: OpenCV not available. Using PIL-based image processing.")
except Exception as e:
    CV2_AVAILABLE = False
    print(f"Warning: OpenCV failed to load: {e}. Using PIL-based image processing.")

# Enhanced text processing
import nltk
from rank_bm25 import BM25Okapi

class EnhancedDocumentProcessor:
    def __init__(self, api_key: str):
        self.api_key = api_key
        self.embeddings = HuggingFaceEmbeddings(
            model_name="sentence-transformers/all-MiniLM-L6-v2",
            model_kwargs={'device': 'cpu'}
        )
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len,
            separators=["\n\n", "\n", " ", ""]
        )
        
        # Enhanced text splitter for different document types
        self.code_splitter = RecursiveCharacterTextSplitter.from_language(
            language="python", chunk_size=800, chunk_overlap=100
        )
        
        # Initialize OCR readers with better error handling
        self.easyocr_reader = None
        self.easyocr_available = EASYOCR_AVAILABLE
        
        if self.easyocr_available:
            try:
                # Initialize EasyOCR with CPU-only mode for cloud compatibility
                self.easyocr_reader = easyocr.Reader(['en'], gpu=False, verbose=False)
                print("✅ EasyOCR initialized successfully")
            except Exception as e:
                print(f"❌ Failed to initialize EasyOCR: {e}")
                self.easyocr_available = False
                self.easyocr_reader = None
            
        # Store availability flags as instance variables
        self.pytesseract_available = PYTESSERACT_AVAILABLE
        self.cv2_available = CV2_AVAILABLE
        
        # Configure Tesseract for better performance (if available)
        if self.pytesseract_available:
            # Enhanced Tesseract configuration for better accuracy
            self.tesseract_config = r'--oem 3 --psm 6'
            # Try different PSM modes for different image types
            self.tesseract_configs = {
                'auto': r'--oem 3 --psm 3',  # Fully automatic page segmentation
                'single_column': r'--oem 3 --psm 4',  # Single column of text
                'single_block': r'--oem 3 --psm 6',  # Single uniform block of text
                'single_line': r'--oem 3 --psm 7',   # Single text line
                'single_word': r'--oem 3 --psm 8',   # Single word
                'sparse': r'--oem 3 --psm 11'        # Sparse text
            }
            print(f"✅ Tesseract configured with {len(self.tesseract_configs)} different recognition modes")
        else:
            self.tesseract_config = None
            self.tesseract_configs = {}
            print("⚠️ Tesseract not available - OCR functionality will be limited")
            
        # Document storage with metadata
        self.documents: List[Document] = []
        self.qa_history = []
        self.document_metadata = {}
        
        # Supported file types
        self.supported_text_types = {
            '.pdf', '.docx', '.txt', '.csv', '.json', '.xml', 
            '.yml', '.yaml', '.xls', '.xlsx', '.pptx', '.html', '.md'
        }
        self.supported_image_types = {
            '.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.tif', '.webp', '.gif'
        }
        
        # Download NLTK data if needed
        try:
            nltk.data.find('tokenizers/punkt')
        except LookupError:
            nltk.download('punkt', quiet=True)
        try:
            nltk.data.find('corpora/stopwords')
        except LookupError:
            nltk.download('stopwords', quiet=True)
        
    def add_qa_to_history(self, question: str, answer: str):
        """Add a Q&A pair to the history."""
        self.qa_history.append({
            'question': question,
            'answer': answer,
            'timestamp': datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        })
    
    def generate_qa_pdf(self, output_path: str, qa_history: List[Dict] = None):
        """Generate a PDF containing the Q&A history."""
        # Use provided history or fall back to internal history
        history_to_use = qa_history if qa_history is not None else self.qa_history
        
        if not history_to_use:
            # Create a simple message if no Q&A history exists
            doc = SimpleDocTemplate(output_path, pagesize=letter)
            styles = getSampleStyleSheet()
            elements = []
            
            elements.append(Paragraph("Document Q&A Session", styles['Title']))
            elements.append(Spacer(1, 20))
            elements.append(Paragraph("No questions and answers recorded in this session.", styles['Normal']))
            elements.append(Spacer(1, 20))
            elements.append(Paragraph(f"Session Date: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}", styles['Normal']))
            
            doc.build(elements)
            return output_path
        
        doc = SimpleDocTemplate(output_path, pagesize=letter)
        styles = getSampleStyleSheet()
        elements = []
        
        # Title
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            spaceAfter=30,
            textColor=colors.darkblue
        )
        elements.append(Paragraph("📚 InsightDocs Q&A Session", title_style))
        
        # Session info
        session_style = ParagraphStyle(
            'SessionInfo',
            parent=styles['Normal'],
            fontSize=12,
            textColor=colors.gray
        )
        elements.append(Paragraph(f"Generated on: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}", session_style))
        elements.append(Paragraph(f"Total Questions: {len(history_to_use)}", session_style))
        elements.append(Spacer(1, 30))
        
        # Q&A pairs
        for i, qa in enumerate(history_to_use, 1):
            # Question
            question_style = ParagraphStyle(
                'Question',
                parent=styles['Heading2'],
                fontSize=14,
                textColor=colors.darkblue,
                spaceAfter=10,
                leftIndent=0
            )
            elements.append(Paragraph(f"Question {i}: {qa['question']}", question_style))
            
            # Answer
            answer_style = ParagraphStyle(
                'Answer',
                parent=styles['Normal'],
                fontSize=11,
                spaceAfter=15,
                leftIndent=20,
                textColor=colors.black
            )
            # Clean up the answer text for PDF
            answer_text = qa['answer'].replace('\n', '<br/>')
            elements.append(Paragraph(f"Answer: {answer_text}", answer_style))
            
            # Timestamp
            timestamp_style = ParagraphStyle(
                'Timestamp',
                parent=styles['Italic'],
                fontSize=9,
                textColor=colors.gray,
                leftIndent=20
            )
            elements.append(Paragraph(f"Asked at: {qa['timestamp']}", timestamp_style))
            elements.append(Spacer(1, 25))
        
        # Footer
        footer_style = ParagraphStyle(
            'Footer',
            parent=styles['Normal'],
            fontSize=10,
            textColor=colors.gray,
            alignment=1  # Center alignment
        )
        elements.append(Spacer(1, 30))
        elements.append(Paragraph("Generated by InsightDocs - Multi-Format RAG Application", footer_style))
        
        # Build PDF
        doc.build(elements)
        return output_path
    
    def read_file(self, file_path: str) -> str:
        """Read different file types and return their content as text."""
        file_extension = os.path.splitext(file_path)[1].lower()
        
        # Check if it's an image file first
        if file_extension in self.supported_image_types:
            return self._read_image(file_path)
        
        # Use python-magic for better file type detection
        if MAGIC_AVAILABLE:
            try:
                mime = magic.Magic(mime=True)
                file_type = mime.from_file(file_path)
            except Exception:
                # Fallback to extension-based detection
                file_type = self._get_mime_from_extension(file_extension)
        else:
            # Use extension-based detection when magic is not available
            file_type = self._get_mime_from_extension(file_extension)
        
        # If we get a generic type, try to determine by extension
        if file_type == 'application/octet-stream' or file_type == 'application/octet-stream':
            file_type = self._get_mime_from_extension(file_extension)
        
        print(f"Processing file: {os.path.basename(file_path)} (extension: {file_extension}, mime: {file_type})")
        
        # Handle different file types
        try:
            if file_type == 'application/pdf' or file_extension == '.pdf':
                return self._read_pdf(file_path)
            elif file_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document' or file_extension == '.docx':
                return self._read_docx(file_path)
            elif file_type == 'text/plain' or file_extension in ['.txt', '.md']:
                return self._read_txt(file_path)
            elif file_type == 'text/csv' or file_extension == '.csv':
                return self._read_csv(file_path)
            elif file_type == 'application/json' or file_extension == '.json':
                return self._read_json(file_path)
            elif file_type == 'text/xml' or file_extension == '.xml':
                return self._read_xml(file_path)
            elif file_type == 'text/yaml' or file_extension in ['.yml', '.yaml']:
                return self._read_yaml(file_path)
            elif file_type == 'application/vnd.ms-excel' or file_extension == '.xls':
                return self._read_xls(file_path)
            elif file_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet' or file_extension == '.xlsx':
                return self._read_xlsx(file_path)
            elif file_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation' or file_extension == '.pptx':
                return self._read_pptx(file_path)
            elif file_type == 'text/html' or file_extension == '.html':
                return self._read_html(file_path)
            else:
                # Try to read as text if file type is unknown but extension suggests text
                if file_extension in ['.txt', '.md', '.log', '.ini', '.cfg', '.conf']:
                    return self._read_txt(file_path)
                else:
                    raise ValueError(f"Unsupported file type: {file_type} for file: {os.path.basename(file_path)}")
        except Exception as e:
            # If specific reader fails, try reading as text as last resort
            if "Unsupported file type" not in str(e):
                try:
                    print(f"Primary reader failed for {os.path.basename(file_path)}, trying as text...")
                    return self._read_txt(file_path)
                except:
                    pass
            raise e
    
    def _get_mime_from_extension(self, extension: str) -> str:
        """Get MIME type from file extension as fallback."""
        mime_map = {
            '.pdf': 'application/pdf',
            '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            '.txt': 'text/plain',
            '.csv': 'text/csv',
            '.json': 'application/json',
            '.xml': 'text/xml',
            '.yml': 'text/yaml',
            '.yaml': 'text/yaml',
            '.xls': 'application/vnd.ms-excel',
            '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            '.html': 'text/html',
            '.md': 'text/plain'
        }
        return mime_map.get(extension, 'application/octet-stream')
    
    def _read_pdf(self, file_path: str) -> str:
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
        return text
    
    def _read_docx(self, file_path: str) -> str:
        doc = Document(file_path)
        return "\n".join([paragraph.text for paragraph in doc.paragraphs])
    
    def _read_txt(self, file_path: str) -> str:
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()
    
    def _read_csv(self, file_path: str) -> str:
        df = pd.read_csv(file_path)
        return df.to_string()
    
    def _read_json(self, file_path: str) -> str:
        with open(file_path, 'r', encoding='utf-8') as file:
            data = json.load(file)
            return json.dumps(data, indent=2)
    
    def _read_xml(self, file_path: str) -> str:
        tree = ET.parse(file_path)
        root = tree.getroot()
        return ET.tostring(root, encoding='unicode', method='xml')
    
    def _read_yaml(self, file_path: str) -> str:
        with open(file_path, 'r', encoding='utf-8') as file:
            data = yaml.safe_load(file)
            return yaml.dump(data, default_flow_style=False)
    
    def _read_xls(self, file_path: str) -> str:
        workbook = xlrd.open_workbook(file_path)
        text = ""
        for sheet in workbook.sheets():
            text += f"Sheet: {sheet.name}\n"
            for row in range(sheet.nrows):
                text += "\t".join(str(sheet.cell_value(row, col)) for col in range(sheet.ncols)) + "\n"
        return text
    
    def _read_xlsx(self, file_path: str) -> str:
        workbook = openpyxl.load_workbook(file_path, data_only=True)
        text = ""
        for sheet in workbook.sheetnames:
            text += f"Sheet: {sheet}\n"
            ws = workbook[sheet]
            for row in ws.rows:
                text += "\t".join(str(cell.value) for cell in row) + "\n"
        return text
    
    def _read_pptx(self, file_path: str) -> str:
        prs = pptx.Presentation(file_path)
        text = ""
        for slide in prs.slides:
            text += f"Slide {slide.slide_number}\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    
    def _read_html(self, file_path: str) -> str:
        with open(file_path, 'r', encoding='utf-8') as file:
            html_content = file.read()
            # Remove HTML tags
            text = re.sub('<[^<]+?>', '', html_content)
            # Decode HTML entities
            text = html.unescape(text)
            return text
    
    def process_documents(self, file_paths: List[str], file_mapping: Dict[str, str] = None) -> Dict[str, Any]:
        """Process multiple documents and create enhanced vector stores."""
        all_chunks = []
        all_documents = []
        file_types = {}
        processing_errors = []
        
        for file_path in file_paths:
            # Get the original filename if mapping is provided
            if file_mapping and file_path in file_mapping:
                original_name = file_mapping[file_path]
                display_name = original_name
            else:
                original_name = os.path.basename(file_path)
                display_name = original_name
            
            try:
                # Read the file content
                content = self.read_file(file_path)
                
                if not content or not content.strip():
                    print(f"Warning: No content extracted from {display_name}")
                    processing_errors.append(f"{display_name}: No readable content found")
                    continue
                
                # Split content into chunks
                file_extension = os.path.splitext(file_path)[1].lower()
                if file_extension in ['.py', '.js', '.ts', '.java', '.cpp', '.c']:
                    chunks = self.code_splitter.split_text(content)
                else:
                    chunks = self.text_splitter.split_text(content)
                
                if not chunks:
                    print(f"Warning: No chunks created from {display_name}")
                    processing_errors.append(f"{display_name}: Could not create text chunks")
                    continue
                
                # Add file metadata to each chunk
                content_type = self._get_content_type(file_extension)
                file_types[display_name] = content_type
                
                # Create documents with metadata
                for i, chunk in enumerate(chunks):
                    if chunk.strip():  # Only add non-empty chunks
                        doc_content = f"[Document: {display_name} | Type: {content_type} | Chunk: {i+1}]\n{chunk}"
                        doc = Document(
                            page_content=doc_content,
                            metadata={
                                "source": display_name,
                                "file_path": file_path,
                                "original_name": original_name,
                                "content_type": content_type,
                                "chunk_id": i+1,
                                "file_extension": file_extension
                            }
                        )
                        all_documents.append(doc)
                        all_chunks.append(doc_content)
                
                print(f"✅ Processed {display_name}: {len([c for c in chunks if c.strip()])} chunks")
                
            except Exception as e:
                error_msg = f"{display_name}: {str(e)}"
                print(f"❌ Error processing {error_msg}")
                processing_errors.append(error_msg)
                continue
        
        # Provide detailed error information
        if not all_chunks:
            error_details = "\n".join(processing_errors) if processing_errors else "Unknown processing error"
            raise ValueError(f"""No text content could be extracted from any of the uploaded files. 

Processing Details:
{error_details}

Troubleshooting Tips:
1. Ensure files are not corrupted or password-protected
2. Check that file extensions match their actual content
3. For images, ensure they contain readable text
4. Try uploading files individually to identify problematic ones
5. Supported formats: PDF, Word, Excel, PowerPoint, Text, Images, etc.""")
        
        # Show processing summary
        if processing_errors:
            print(f"\n⚠️ Processing Summary:")
            print(f"✅ Successfully processed: {len(file_types)} files")
            print(f"❌ Failed to process: {len(processing_errors)} files")
            for error in processing_errors:
                print(f"   - {error}")
        
        # Create FAISS vector store
        vector_store = FAISS.from_documents(all_documents, self.embeddings)
        
        # Create BM25 retriever for keyword-based search
        bm25_retriever = BM25Retriever.from_documents(all_documents)
        bm25_retriever.k = 5
        
        # Store for later use in ensemble retrieval
        self.documents = all_documents
        self.document_metadata = {
            'file_types': file_types,
            'total_chunks': len(all_chunks),
            'total_files': len(file_types),
            'processing_errors': processing_errors
        }
        
        return {
            'vector_store': vector_store,
            'bm25_retriever': bm25_retriever,
            'documents': all_documents,
            'metadata': self.document_metadata
        }
    
    def _get_content_type(self, file_extension: str) -> str:
        """Get a human-readable content type based on file extension."""
        content_types = {
            '.pdf': 'PDF Document',
            '.docx': 'Word Document',
            '.txt': 'Text File',
            '.csv': 'CSV Spreadsheet',
            '.json': 'JSON Data',
            '.xml': 'XML Document',
            '.yml': 'YAML File',
            '.yaml': 'YAML File',
            '.xls': 'Excel Spreadsheet',
            '.xlsx': 'Excel Spreadsheet',
            '.pptx': 'PowerPoint Presentation',
            '.html': 'HTML Document',
            '.md': 'Markdown Document',
            '.jpg': 'JPEG Image',
            '.jpeg': 'JPEG Image',
            '.png': 'PNG Image',
            '.bmp': 'BMP Image',
            '.tiff': 'TIFF Image',
            '.tif': 'TIFF Image',
            '.webp': 'WebP Image',
            '.gif': 'GIF Image',
            '.py': 'Python Code',
            '.js': 'JavaScript Code',
            '.ts': 'TypeScript Code',
            '.java': 'Java Code',
            '.cpp': 'C++ Code',
            '.c': 'C Code'
        }
        return content_types.get(file_extension, 'Unknown Type')
    
    def create_qa_chain(self, document_store: Dict[str, Any]) -> RetrievalQA:
        """Create an enhanced question-answering chain using ensemble retrieval."""
        llm = ChatGroq(
            groq_api_key=self.api_key,
            model_name="gemma2-9b-it",
            temperature=0.3,  # Lower temperature for more focused answers
            max_tokens=1024  # Adjusted to stay within model limits
        )
        
        # Get retrievers from the document store
        vector_store = document_store['vector_store']
        bm25_retriever = document_store['bm25_retriever']
        
        # Create ensemble retriever combining semantic and keyword search
        vector_retriever = vector_store.as_retriever(
            search_kwargs={
                "k": 8,  # Increased for better context
                "fetch_k": 20,
                "lambda_mult": 0.6  # Balance between relevance and diversity
            }
        )
        
        # Combine retrievers with different weights
        ensemble_retriever = EnsembleRetriever(
            retrievers=[vector_retriever, bm25_retriever],
            weights=[0.6, 0.4]  # Favor semantic search slightly
        )
        
        # Create QA chain with custom prompt
        qa_chain = RetrievalQA.from_chain_type(
            llm=llm,
            chain_type="stuff",
            retriever=ensemble_retriever,
            return_source_documents=True,
            chain_type_kwargs={
                "prompt": self._create_custom_prompt()
            }
        )
        return qa_chain
    
    def _create_custom_prompt(self):
        """Create a custom prompt template for better responses."""
        from langchain.prompts import PromptTemplate
        
        template = """Use the following pieces of context to answer the question at the end. 
        The context may include text from various document types including PDFs, Word documents, spreadsheets, presentations, images with OCR text, and other file formats.
        
        If you don't know the answer, just say that you don't know, don't try to make up an answer.
        When answering, consider the document type and source of the information.
        
        Context:
        {context}
        
        Question: {question}
        
        Helpful Answer:"""
        
        return PromptTemplate(
            template=template,
            input_variables=["context", "question"]
        )
    
    def preprocess_image_pil(self, image_path: str) -> Image.Image:
        """Preprocess image using PIL for better OCR accuracy when OpenCV is not available."""
        try:
            # Open and convert to RGB if needed
            img = Image.open(image_path)
            print(f"📸 Original image: {img.size} pixels, mode: {img.mode}")
            
            # Convert to RGB first if needed
            if img.mode in ('RGBA', 'P'):
                background = Image.new('RGB', img.size, (255, 255, 255))
                if img.mode == 'P':
                    img = img.convert('RGBA')
                background.paste(img, mask=img.split()[-1] if img.mode == 'RGBA' else None)
                img = background
            elif img.mode not in ('RGB', 'L'):
                img = img.convert('RGB')
            
            # Convert to grayscale for better OCR
            if img.mode != 'L':
                img = img.convert('L')
            
            # Resize if image is too small (OCR works better with larger text)
            width, height = img.size
            min_size = 800  # Minimum size for good OCR
            if width < min_size or height < min_size:
                scale_factor = max(min_size/width, min_size/height)
                new_size = (int(width * scale_factor), int(height * scale_factor))
                img = img.resize(new_size, Image.Resampling.LANCZOS)
                print(f"📏 Resized image to: {new_size}")
            
            # Resize if image is too large (can cause memory issues)
            max_size = 3000
            if width > max_size or height > max_size:
                scale_factor = min(max_size/width, max_size/height)
                new_size = (int(width * scale_factor), int(height * scale_factor))
                img = img.resize(new_size, Image.Resampling.LANCZOS)
                print(f"📏 Reduced image size to: {new_size}")
            
            # Enhance contrast
            enhancer = ImageEnhance.Contrast(img)
            img = enhancer.enhance(1.5)  # More conservative enhancement
            
            # Enhance sharpness
            enhancer = ImageEnhance.Sharpness(img)
            img = enhancer.enhance(1.3)  # More conservative enhancement
            
            # Apply filters for noise reduction
            img = img.filter(ImageFilter.MedianFilter(size=3))
            
            print(f"✅ PIL preprocessing complete: {img.size}")
            return img
            
        except Exception as e:
            print(f"❌ PIL preprocessing failed: {e}")
            # Return original image if preprocessing fails
            return Image.open(image_path)

    def preprocess_image(self, image_path: str) -> np.ndarray:
        """Preprocess image for better OCR accuracy with fallback options."""
        if not self.cv2_available:
            # Use PIL-based preprocessing
            try:
                pil_img = self.preprocess_image_pil(image_path)
                return np.array(pil_img)
            except Exception as e:
                print(f"PIL preprocessing failed: {e}")
                # Final fallback - just convert to grayscale
                pil_img = Image.open(image_path).convert('L')
                return np.array(pil_img)
        
        try:
            # Read image with OpenCV
            img = cv2.imread(image_path)
            if img is None:
                # Try with PIL if OpenCV fails
                pil_img = Image.open(image_path)
                img = cv2.cvtColor(np.array(pil_img), cv2.COLOR_RGB2BGR)
            
            # Convert to grayscale
            gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
            
            # Apply denoising
            denoised = cv2.fastNlMeansDenoising(gray)
            
            # Apply adaptive thresholding
            thresh = cv2.adaptiveThreshold(
                denoised, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY, 11, 2
            )
            
            # Morphological operations to clean up
            kernel = np.ones((1, 1), np.uint8)
            cleaned = cv2.morphologyEx(thresh, cv2.MORPH_CLOSE, kernel)
            
            return cleaned
        except Exception as e:
            print(f"OpenCV preprocessing failed: {e}, falling back to PIL")
            # Fallback to PIL processing
            return self.preprocess_image_pil(image_path)
    
    def extract_text_from_image(self, image_path: str) -> str:
        """Extract text from image using multiple OCR methods for maximum accuracy."""
        print(f"🔍 Analyzing image: {os.path.basename(image_path)}")
        extracted_texts = []
        
        # Method 1: EasyOCR (if available)
        if self.easyocr_available and self.easyocr_reader:
            try:
                print("📖 Running EasyOCR analysis...")
                results = self.easyocr_reader.readtext(image_path, detail=1, paragraph=False)
                easyocr_text = []
                confidence_scores = []
                
                for (bbox, text, confidence) in results:
                    if confidence > 0.1:  # Very low threshold to catch everything
                        easyocr_text.append(text.strip())
                        confidence_scores.append(confidence)
                
                if easyocr_text:
                    combined_text = " ".join(easyocr_text)
                    avg_confidence = sum(confidence_scores) / len(confidence_scores)
                    extracted_texts.append(("EasyOCR", combined_text, avg_confidence))
                    print(f"✅ EasyOCR extracted {len(easyocr_text)} text segments (avg confidence: {avg_confidence:.2f})")
                else:
                    print("⚠️ EasyOCR found no text")
                    
            except Exception as e:
                print(f"❌ EasyOCR failed: {e}")
        else:
            print("⚠️ EasyOCR not available")
        
        # Method 2: Tesseract with multiple configurations (if available)
        if self.pytesseract_available:
            print(f"📖 Running Tesseract with {len(self.tesseract_configs)} configurations...")
            
            for config_name, config in self.tesseract_configs.items():
                try:
                    # Try with preprocessed image first
                    try:
                        processed_img = self.preprocess_image_pil(image_path)
                        tesseract_text = pytesseract.image_to_string(processed_img, config=config).strip()
                        
                        if tesseract_text and len(tesseract_text) > 1:
                            # Calculate confidence based on text characteristics
                            confidence = self._estimate_text_quality(tesseract_text)
                            extracted_texts.append((f"Tesseract_{config_name}_processed", tesseract_text, confidence))
                            print(f"✅ Tesseract {config_name} (processed): {len(tesseract_text)} chars (confidence: {confidence:.2f})")
                    except Exception as e:
                        print(f"⚠️ Tesseract {config_name} with preprocessing failed: {e}")
                    
                    # Try with original image as fallback
                    try:
                        original_img = Image.open(image_path)
                        original_text = pytesseract.image_to_string(original_img, config=config).strip()
                        
                        if original_text and len(original_text) > 1:
                            confidence = self._estimate_text_quality(original_text)
                            extracted_texts.append((f"Tesseract_{config_name}_original", original_text, confidence))
                            print(f"✅ Tesseract {config_name} (original): {len(original_text)} chars (confidence: {confidence:.2f})")
                    except Exception as e:
                        print(f"⚠️ Tesseract {config_name} on original failed: {e}")
                        
                except Exception as e:
                    print(f"❌ Tesseract {config_name} configuration failed: {e}")
        else:
            print("⚠️ Tesseract not available")
        
        # Method 3: Try basic PIL-only approach (fallback)
        if not extracted_texts:
            try:
                print("🔄 Trying basic PIL extraction as last resort...")
                img = Image.open(image_path)
                # Convert to text-friendly format
                if img.mode != 'L':
                    img = img.convert('L')
                
                # Try with basic Tesseract if available
                if self.pytesseract_available:
                    basic_text = pytesseract.image_to_string(img).strip()
                    if basic_text:
                        confidence = self._estimate_text_quality(basic_text)
                        extracted_texts.append(("Basic_Tesseract", basic_text, confidence))
                        print(f"✅ Basic extraction: {len(basic_text)} chars")
            except Exception as e:
                print(f"❌ Basic extraction failed: {e}")
        
        # Analyze results
        if not extracted_texts:
            return """❌ TEXT EXTRACTION FAILED

🔍 ANALYSIS SUMMARY:
• No readable text could be detected in this image
• Multiple OCR methods were attempted but found no text content

🔧 POSSIBLE REASONS:
• The image may not contain any text
• Text might be too small, blurry, or distorted
• Text color may be too similar to background
• The image might be a pure graphic/photo without text
• OCR libraries may not be properly configured

💡 SUGGESTIONS:
• Ensure the image contains clear, readable text
• Try images with high contrast between text and background
• Use images with text size of at least 12pt
• Avoid heavily compressed or low-resolution images"""
        
        # Sort by confidence and length
        ranked_texts = sorted(extracted_texts, key=lambda x: (x[2], len(x[1])), reverse=True)
        
        # Get the best result
        best_method, best_text, best_confidence = ranked_texts[0]
        
        # Clean up the text
        best_text = self._clean_extracted_text(best_text)
        
        # If we have multiple good results, try to combine unique content
        if len(ranked_texts) > 1:
            print(f"🔄 Analyzing {len(ranked_texts)} OCR results for combination...")
            
            unique_additions = []
            for method, text, confidence in ranked_texts[1:]:
                if confidence > 0.3:  # Only consider reasonably confident alternatives
                    cleaned_text = self._clean_extracted_text(text)
                    if len(cleaned_text) > len(best_text) * 0.1:  # At least 10% of main text
                        unique_additions.append(f"\n[Alternative OCR result from {method}]:\n{cleaned_text}")
            
            if unique_additions:
                best_text += "\n" + "\n".join(unique_additions[:2])  # Limit to 2 alternatives
        
        print(f"🎯 Best result from {best_method} (confidence: {best_confidence:.2f})")
        print(f"📝 Final extracted text: {len(best_text)} characters, {len(best_text.split())} words")
        
        return best_text

    def _estimate_text_quality(self, text: str) -> float:
        """Estimate the quality/confidence of extracted text."""
        if not text or len(text) < 2:
            return 0.0
        
        # Basic quality indicators
        char_variety = len(set(text.lower())) / len(text) if text else 0
        length_score = min(len(text) / 100, 1.0)
        
        # Check for reasonable text patterns
        letter_count = sum(1 for c in text if c.isalpha())
        letter_ratio = letter_count / len(text) if text else 0
        
        # Penalize text with too many special characters or numbers only
        alpha_ratio = letter_ratio
        if alpha_ratio < 0.3:
            alpha_ratio *= 0.5
        
        # Combine factors
        confidence = (char_variety * 0.3 + length_score * 0.3 + alpha_ratio * 0.4)
        return min(confidence, 1.0)

    def _clean_extracted_text(self, text: str) -> str:
        """Clean and normalize extracted text."""
        if not text:
            return ""
        
        # Remove excessive whitespace
        text = re.sub(r'\s+', ' ', text)
        
        # Remove leading/trailing whitespace
        text = text.strip()
        
        # Remove weird characters that OCR sometimes produces
        text = re.sub(r'[^\w\s\-.,;:!?()\[\]{}"\'/\\@#$%^&*+=<>|`~]', '', text)
        
        return text
    
    def analyze_image_metadata(self, image_path: str) -> Dict[str, Any]:
        """Extract comprehensive metadata from image."""
        try:
            with Image.open(image_path) as img:
                metadata = {
                    'filename': os.path.basename(image_path),
                    'format': img.format,
                    'mode': img.mode,
                    'size': img.size,
                    'width': img.width,
                    'height': img.height,
                    'file_size': os.path.getsize(image_path)
                }
                
                # Calculate image characteristics
                metadata['aspect_ratio'] = round(img.width / img.height, 2)
                metadata['megapixels'] = round((img.width * img.height) / 1000000, 2)
                
                # Analyze image complexity (helps determine OCR approach)
                if img.mode in ('RGB', 'RGBA'):
                    # Convert to grayscale to analyze contrast
                    gray_img = img.convert('L')
                    img_array = np.array(gray_img)
                    metadata['mean_brightness'] = round(np.mean(img_array), 2)
                    metadata['contrast_std'] = round(np.std(img_array), 2)
                    
                    # Estimate text likelihood based on contrast
                    if metadata['contrast_std'] > 50:
                        metadata['text_likelihood'] = 'High'
                    elif metadata['contrast_std'] > 25:
                        metadata['text_likelihood'] = 'Medium'
                    else:
                        metadata['text_likelihood'] = 'Low'
                
                # Extract EXIF data if available
                try:
                    exif_data = img._getexif()
                    if exif_data:
                        metadata['has_exif'] = True
                        # Extract useful EXIF tags
                        useful_tags = {
                            'DateTime': 306,
                            'Software': 305,
                            'ImageDescription': 270,
                            'Make': 271,
                            'Model': 272
                        }
                        
                        for tag_name, tag_id in useful_tags.items():
                            if tag_id in exif_data:
                                metadata[f'exif_{tag_name.lower()}'] = str(exif_data[tag_id])
                    else:
                        metadata['has_exif'] = False
                except:
                    metadata['has_exif'] = False
                
                return metadata
        except Exception as e:
            return {
                'filename': os.path.basename(image_path),
                'error': f"Could not analyze image metadata: {str(e)}"
            }

    def _read_image(self, file_path: str) -> str:
        """Process image file and extract comprehensive text content."""
        print(f"🖼️ Processing image: {os.path.basename(file_path)}")
        
        try:
            # Get image metadata first
            metadata = self.analyze_image_metadata(file_path)
            
            # Extract text using enhanced OCR
            text_content = self.extract_text_from_image(file_path)
            
            # Determine content quality and provide insights
            word_count = len(text_content.split()) if text_content and not text_content.startswith("❌") else 0
            char_count = len(text_content) if text_content and not text_content.startswith("❌") else 0
            
            # Quality assessment
            quality_indicators = []
            extraction_success = True
            
            if text_content.startswith("❌"):
                quality_indicators.append("❌ No text detected")
                extraction_success = False
            elif word_count > 100:
                quality_indicators.append("✅ Rich text content")
            elif word_count > 20:
                quality_indicators.append("✅ Moderate text content")
            elif word_count > 5:
                quality_indicators.append("⚠️ Limited text content")
            elif word_count > 0:
                quality_indicators.append("⚠️ Minimal text detected")
            else:
                quality_indicators.append("❌ No readable text found")
                extraction_success = False
            
            # Image characteristics that affect OCR
            if metadata.get('text_likelihood') == 'High':
                quality_indicators.append("✅ High contrast (good for OCR)")
            elif metadata.get('text_likelihood') == 'Medium':
                quality_indicators.append("⚠️ Medium contrast")
            elif metadata.get('text_likelihood') == 'Low':
                quality_indicators.append("❌ Low contrast (poor for OCR)")
            
            # Size assessment
            width = metadata.get('width', 0)
            height = metadata.get('height', 0)
            if width > 1000 and height > 1000:
                quality_indicators.append("✅ Good resolution for OCR")
            elif width > 500 and height > 500:
                quality_indicators.append("⚠️ Adequate resolution")
            else:
                quality_indicators.append("❌ Low resolution (may affect OCR)")
            
            # Create comprehensive formatted content
            content = f"""📷 IMAGE ANALYSIS REPORT
════════════════════════════════════════════════════════════════════════

📋 FILE INFORMATION:
• Filename: {metadata.get('filename', 'Unknown')}
• Format: {metadata.get('format', 'Unknown')}
• Dimensions: {width:,} × {height:,} pixels ({metadata.get('megapixels', 0):.1f} MP)
• File Size: {metadata.get('file_size', 0):,} bytes
• Color Mode: {metadata.get('mode', 'Unknown')}
• Aspect Ratio: {metadata.get('aspect_ratio', 'Unknown')}

📊 OCR ANALYSIS:
• Text Extraction: {'✅ SUCCESS' if extraction_success else '❌ FAILED'}
• Characters Detected: {char_count:,}
• Words Detected: {word_count:,}
• Quality Indicators: {' | '.join(quality_indicators)}

� IMAGE CHARACTERISTICS:
• Brightness Level: {metadata.get('mean_brightness', 'Unknown')}
• Contrast Score: {metadata.get('contrast_std', 'Unknown')}
• OCR Likelihood: {metadata.get('text_likelihood', 'Unknown')}

{f"📅 Creation Date: {metadata.get('exif_datetime', 'Not available')}" if metadata.get('exif_datetime') else ""}
{f"📱 Device Info: {metadata.get('exif_make', '')} {metadata.get('exif_model', '')}".strip() if metadata.get('exif_make') else ""}

════════════════════════════════════════════════════════════════════════
📝 EXTRACTED TEXT CONTENT:

{text_content}

════════════════════════════════════════════════════════════════════════
🏷️ DOCUMENT TAGS: Image Document, OCR Processed, Visual Content, {'Text Extracted' if extraction_success else 'No Text Found'}

💡 OCR PERFORMANCE NOTES:
{self._generate_ocr_performance_notes(metadata, extraction_success, word_count)}
"""
            
            print(f"✅ Image processing complete: {word_count} words extracted from {os.path.basename(file_path)}")
            return content
            
        except Exception as e:
            error_content = f"""❌ IMAGE PROCESSING ERROR
════════════════════════════════════════════════════════════════════════

📋 FILE INFORMATION:
• Filename: {os.path.basename(file_path)}
• Status: ❌ Processing Failed
• Error Type: {type(e).__name__}

❌ ERROR DETAILS:
{str(e)}

🔧 TROUBLESHOOTING SUGGESTIONS:
• Verify the image file is not corrupted
• Ensure the image format is supported (JPG, PNG, TIFF, BMP, etc.)
• Check if the image file size is reasonable (not too large/small)
• Try converting the image to PNG or JPEG format
• Ensure OCR libraries are properly installed

════════════════════════════════════════════════════════════════════════
🏷️ DOCUMENT TAGS: Image Document, Processing Error, Failed OCR
"""
            print(f"❌ Error processing image {os.path.basename(file_path)}: {str(e)}")
            return error_content

    def _generate_ocr_performance_notes(self, metadata: dict, success: bool, word_count: int) -> str:
        """Generate performance notes based on OCR results and image characteristics."""
        notes = []
        
        if success and word_count > 50:
            notes.append("✅ Excellent OCR performance - image has clear, readable text")
        elif success and word_count > 10:
            notes.append("✅ Good OCR performance - successfully extracted readable text")
        elif success and word_count > 0:
            notes.append("⚠️ Limited OCR results - image may have small or unclear text")
        else:
            notes.append("❌ OCR extraction failed - see suggestions below")
        
        # Image quality notes
        contrast = metadata.get('contrast_std', 0)
        if isinstance(contrast, (int, float)):
            if contrast > 60:
                notes.append("✅ High contrast image - ideal for OCR")
            elif contrast > 30:
                notes.append("⚠️ Medium contrast - acceptable for OCR")
            else:
                notes.append("❌ Low contrast - may hinder text recognition")
        
        # Size recommendations
        width = metadata.get('width', 0)
        height = metadata.get('height', 0)
        if width < 500 or height < 500:
            notes.append("💡 TIP: Higher resolution images (>1000px) typically yield better OCR results")
        
        if not success:
            notes.extend([
                "💡 TIP: Ensure text is dark on light background (or vice versa)",
                "💡 TIP: Avoid images with decorative fonts or stylized text",
                "💡 TIP: Text should be horizontally oriented for best results"
            ])
        
        return "\n".join(f"• {note}" for note in notes)

# Backward compatibility alias
DocumentProcessor = EnhancedDocumentProcessor